from flask import Flask, request, jsonify

app = Flask(__name__)


@app.route('/run-script', methods=['POST'])
def run_script():
    # Retrieve data from the request
    data = request.get_json()

    # Your Python script logic here

    # libraries

    from office365.runtime.auth.authentication_context import AuthenticationContext
    from office365.sharepoint.client_context import ClientContext
    from office365.runtime.auth.user_credential import UserCredential
    from office365.sharepoint.files.file import File
    from office365.sharepoint.files.creation_information import FileCreationInformation

    import pptx
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.dml import MSO_THEME_COLOR
    import collections
    import collections.abc

    from urllib import response
    from urllib.parse import parse_qs, urlparse, quote
    import urllib.parse

    import sys
    import os
    import time
    import pendulum
    import datetime
    from io import StringIO, BytesIO

    import json
    import pandas as pd
    import polars as pl
    import numpy as np
    import matplotlib as mpl

    ######################################################################################################################

    today_date = pendulum.today().date()
    print(today_date)

    USERNAME = "ar_serviceaccount@alpharecon.com"
    PASSWORD = "h9XcwVW8E5hD"

    ######################################################################################################################

    # create Sharepoint class object

    class SharePointClient:
        def __init__(self, site_url, username, password):
            self.site_url = site_url
            self.username = username
            self.password = password
            self.client_context = self._get_client_context()

        def _get_client_context(self):
            ctx_auth = AuthenticationContext(url=self.site_url)
            if ctx_auth.acquire_token_for_user(username=self.username, password=self.password):
                client_ctx = ClientContext(self.site_url, ctx_auth)
                return client_ctx
            else:
                raise ValueError("Authentication failed")

        def view_files(self, library_path, type):

            if type == "file":
                library = self.client_context.web.get_folder_by_server_relative_url(
                    library_path)
                files = library.files
                self.client_context.load(files)
                self.client_context.execute_query()

                undeleted_files = [{'file_name': f.properties['Name'], 'file_path': f.properties['ServerRelativeUrl']} for f in files if not f.listItemAllFields.is_property_available("EffectiveBasePermissions") or
                                   not f.listItemAllFields.effective_base_permissions.high & 0x1]
                return undeleted_files

            if type == "folder":
                folder = self.client_context.web.get_folder_by_server_relative_url(
                    library_path)
                folders = folder.folders
                self.client_context.load(folders)
                self.client_context.execute_query()

                if not folders:
                    print("No folders found.")
                    return

                undeleted_folders = [{'folder_name': f.properties['Name'], 'folder_path': f.properties['ServerRelativeUrl']} for f in folders if not f.list_item_all_fields.is_property_available("EffectiveBasePermissions") or
                                     not f.list_item_all_fields.effective_base_permissions.high & 0x1]
                return undeleted_folders

            else:
                print(
                    "Please identify if you want to view a 'file' or 'folder' in type.")

        def upload_file(self, file_content, destination_folder_url, file_name, overwrite=False):
            folder = self.client_context.web.get_folder_by_server_relative_url(
                destination_folder_url)
            files = folder.files
            self.client_context.load(files)
            self.client_context.execute_query()

            # Use ServerRelativeUrl to create the file path in the destination folder
            # file_name = local_path.split("/")[-1]
            file_path = f"{destination_folder_url}/{file_name}"

            print('upload file_path', file_path)

            # Check if the file already exists
            existing_file = next(
                (f for f in files if f.properties['Name'] == file_name), None)

            # if existing_file:
            #     if overwrite:
            #         # Delete the existing file if overwrite is True
            #         existing_file.delete_object()
            #         self.client_context.execute_query()

            if existing_file:
                if overwrite:
                    # Delete the existing file if overwrite is True
                    existing_file.delete_object()
                    self.client_context.execute_query()
                else:
                    # If not overwriting and file exists, raise an exception or handle accordingly
                    raise ValueError(f"File '{file_path}' already exists.")

            # Upload the new file
            with BytesIO(file_content.getvalue()) as file_info:
                # Use the FileCollection.add method to add the file to the folder
                uploaded_file = folder.files.add(
                    file_name, file_info.read(), overwrite=overwrite)
                self.client_context.execute_query()

            print(
                f"File uploaded successfully. Server relative url: {file_path}")

        def download_upload_file(self, file_name, sharepoint_path_get, sharepoint_path_push, data):
            folder = self.client_context.web.get_folder_by_server_relative_url(
                sharepoint_path_get)
            files = folder.files
            self.client_context.load(files)
            self.client_context.execute_query()

            # Find the file by name
            # target_file = next(
            #     (f for f in files if f.properties['Name'] == file_name), None)
            for f in files:
                if f.properties['Name'] == file_name:
                    target_file = f

            if target_file:
                # Download the file content
                file_content = target_file.read()
                self.client_context.execute_query()

                prs = pptx.Presentation(BytesIO(file_content))

                # Define a function to replace placeholders
                def replace_text_runs(runs, data):
                    for run in runs:
                        for k, v in data.items():
                            # Using replace method to replace all occurrences of the placeholder
                            run.text = run.text.replace(
                                str(k), str(v) if v is not None else '')

                # Access slides and modify content as needed
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for paragraph in shape.text_frame.paragraphs:
                                replace_text_runs(paragraph.runs, data)

                # Save modified presentation with a new name
                output_file_name = f"{data['$$1$$']}_{data['$$44$$']}_RECON_REQUEST.pptx"
                updated_content = BytesIO()
                prs.save(updated_content)

                # Upload the modified content
                self.upload_file(file_content=updated_content, destination_folder_url=sharepoint_path_push,
                                 file_name=output_file_name, overwrite=False)
            else:
                print(f"File '{file_name}' not found in the specified folder.")

        def read_txt_file_content(self, file_path):
            sp_file = self.client_context.web.get_file_by_server_relative_path(
                file_path)
            self.client_context.load(sp_file)
            self.client_context.execute_query()

            if sp_file.properties["ServerRelativeUrl"].endswith(".txt"):
                file_content = sp_file.read()
                decoded_content = file_content.decode('utf-8')
                return decoded_content
            else:
                print(f"'{file_path}' is not a .txt file")

            # sp_file = self.client_context.web.get_file_by_url(
            #     file_path)
            # self.client_context.load(sp_file)
            # self.client_context.execute_query()

            # if sp_file.properties["ServerRelativeUrl"].endswith(".txt"):
            #     file_content = sp_file.read()
            #     decoded_content = file_content.decode('utf-8')
            #     return decoded_content
            # else:
            #     print(f"'{file_path}' is not a .txt file")

    ######################################################################################################################

    # Sharepoint url for requests
    request_site_url = "https://alpharecon0.sharepoint.com/sites/IntelProductSMT"
    # request_folder_url = "/sites/IntelProductSMT/Shared%20Documents/Document%20Repository/Intel%20Ops/Falcon%20Test%20Recons/"
    request_folder_url = "/sites/IntelProductSMT/Shared%20Documents/Document%20Repository/Intel%20Ops/Automated%20Recon%20Reports/"

    # Authenticate sharepoint client
    sharepoint_client_request = SharePointClient(
        request_site_url, USERNAME, PASSWORD)

    # Look for Output and Report Template folders within request folder
    folders = sharepoint_client_request.view_files(
        request_folder_url, 'folder')

    # Get the ServerRelativeUrl path to folders neeeded
    for i, v in enumerate(folders):
        if v['folder_name'] == 'Output':
            library_path_to_upload = v['folder_path']
        if v['folder_name'] == 'Report Templates':
            library_path_to_download = v['folder_path']
        if v['folder_name'] == 'Request Details':
            detailed_requests = v['folder_path']

    report_template_files_to_read = sharepoint_client_request.view_files(
        library_path_to_download, 'file')

    # View files in the specified library
    request_files_to_read = sharepoint_client_request.view_files(
        detailed_requests, 'file')
    ######################################################################################################################
    # Report Creation

    # Files in request folder to store into data dictionary
    for i, v in enumerate(request_files_to_read):
        print('TESTING', v['file_name'])
        if "INTSUM" in v['file_name']:
            report_type = "INTSUM"
            print(f"Creating {report_type} report from {v['file_name']}.")
        elif "TVA" in v['file_name']:
            report_type = "TVA"
            print(f"Creating {report_type} report from {v['file_name']}.")
        elif "SITREP" in v['file_name']:
            report_type = "SITREP"
            print(f"Creating {report_type} report from {v['file_name']}.")
        else:
            print(
                f"Report type could not be concluded based on {v['file_name']}")

        # Read in data from .txt requests
        try:
            text_info = sharepoint_client_request.read_txt_file_content(
                v['file_path'])

        except:
            f"Cannot get content for {v['file_name']}"

        # Store in dictionary
        try:
            company = text_info.split(" -")[0]
            print(company)
            report_data = {
                '$$1$$': company.strip(), '$$44$$': report_type.strip()}
        except:
            "No information"

        ######################################################################################################################

        for k, y in enumerate(report_template_files_to_read):
            if report_type in y['file_name']:

                # Download a template from SharePoint based on Report name
                download_file_name = y['file_name']

            else:
                print('Not enough information to discern which report template.')

        # download the report template, write data to placeholder variables and upload into the output folder
        sharepoint_client_request.download_upload_file(
            download_file_name, library_path_to_download, library_path_to_upload, data=report_data)

        ######################################################################################################################
        print(f'{report_type} report generated for {company}.')
        print('='*100)

    print('All reports generated from requests in folder!')

    result = {'message': 'Python script executed successfully', 'data': data}

    return jsonify(result)


if __name__ == '__main__':
    app.run(debug=True)
